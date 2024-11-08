from bs4 import BeautifulSoup
import requests
from pptx import Presentation
from pptx.util import Inches
import os

# Function to add images to a slide
def add_images_to_slide(prs, images, slide_title):
    # Create a blank slide layout
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title to the slide
    title_placeholder = slide.shapes.title
    title_placeholder.text = slide_title

    # Define the position for the first image
    left = Inches(1)
    top = Inches(1)
    img_width = Inches(3)
    img_height = Inches(3)

    # Loop through images and add to the slide
    for i, img_path in enumerate(images):
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, img_width, img_height)
            left += img_width  # Move the position for the next image to the right
        else:
            print(f"Image not found: {img_path}")
    
    return prs

# Open and parse the HTML file
with open('cnvAnalysis_v2.html', 'r', encoding='utf-8') as file:
    soup = BeautifulSoup(file, 'html.parser')

# Find the Tumors section by its id
tumors_tab = soup.find('div', {'id': 'tumors', 'class': 'section level2 tabset'})

if tumors_tab is None:
    print("Tumors tab not found. Check the HTML structure.")
else:
    # Find all direct div children of the Tumors tab that are sub-tabs (don't rely on specific class)
    sub_tabs = tumors_tab.find_all('div', recursive=False)

    if not sub_tabs:
        print("No sub-tabs found. Check the HTML structure.")
    else:
        # Create a PowerPoint presentation
        prs = Presentation()

        # Create a folder to store downloaded images
        os.makedirs('downloaded_images', exist_ok=True)

        # Loop through each sub-tab
        for sub_tab in sub_tabs:
            # Extract the sub-tab title (e.g., "352")
            sub_tab_title = sub_tab.find('h3').text if sub_tab.find('h3') else "Unnamed Sub-tab"

            # Find all images in the current sub-tab
            images = sub_tab.find_all('img')

            if not images:
                print(f"No images found in sub-tab {sub_tab_title}")
                continue

            image_paths = []

            # Loop through each image, download, and save
            for i, img in enumerate(images):
                img_url = img['src']  # Adjust if images are not absolute URLs

                # If the image URL is relative, prepend the base URL
                if not img_url.startswith('http'):
                    img_url = f'https://your_base_url_here/{img_url}'  # Adjust as needed

                try:
                    # Download image
                    img_data = requests.get(img_url).content
                    img_path = f'downloaded_images/{sub_tab_title}_image_{i}.jpg'
                    with open(img_path, 'wb') as img_file:
                        img_file.write(img_data)
                    print(f"Downloaded image {img_path}")
                    image_paths.append(img_path)
                except Exception as e:
                    print(f"Failed to download image from {img_url}: {e}")
                    continue

            # Add images from the current sub-tab to one PowerPoint slide
            prs = add_images_to_slide(prs, image_paths, sub_tab_title)

        # Save the PowerPoint file
        prs.save('output_presentation.pptx')
        print("Images extracted and PowerPoint created.")
